"""
Investor slide deck generation for a ReportSpec, using python-pptx.

Save as: board/extraction/report_deck.py

Deliberately not a 5th Gemini pipeline: every number on a "stat slide"
comes straight from figures already computed elsewhere in this codebase
(commentary.py's SUMMARY_BUILDERS, FinancialPeriod's cross-statement
properties) — zero fabrication risk since nothing here is generated
text. The only AI-authored content that can appear is 1-2 bullet lines
per section pulled from ReportSpec.tailored_narrative, and only when
narrative_approved is True (see ReportSpecViewSet.generate_deck's
fallback rule in views.py).

Pure Python, no native binary dependencies — deliberately lighter-weight
than the Playwright PDF path (report_pdf.py) rather than piling a
second browser/binary dependency onto this app.
"""

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .commentary import SUMMARY_BUILDERS
from .report_narrative import SECTION_LABELS

# Matches tokens.css's --color-forest / --color-ink / --color-grey-text
FOREST = RGBColor(0x34, 0x57, 0xE8)
INK = RGBColor(0x14, 0x18, 0x1F)
GREY = RGBColor(0x6B, 0x72, 0x80)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# (label, summary-dict key or FinancialPeriod attribute, format kind).
# outlook pulls straight off FinancialPeriod's own cross-statement
# properties (roce_pct/dscr/yoy_revenue_growth_pct) rather than
# _outlook_summary's dict, since that dict prefixes borrowed keys
# (e.g. "revenue_growth.revenue") for prompt-building — not what a
# stat slide wants.
STAT_SLIDE_FIELDS = {
    "revenue_growth": [("Revenue", "revenue", "eur"), ("Gross Margin", "gross_margin_pct", "pct")],
    "profitability": [
        ("EBITDA", "ebitda", "eur"),
        ("EBITDA Margin", "ebitda_margin_pct", "pct"),
        ("Operating Loss", "operating_loss", "eur"),
    ],
    "cash_liquidity": [
        ("Cash on Hand", "closing_cash", "eur"),
        ("Cash Runway", "cash_runway_months", "months"),
        ("Free Cash Flow", "free_cash_flow", "eur"),
    ],
    "solvency_leverage": [("Net Assets", "net_assets", "eur"), ("Total Liabilities", "total_liabilities", "eur")],
    "returns": [
        ("Market Cap", "market_cap", "eur"),
        ("Total Customers", "total_customers", "count"),
        ("Revenue per Customer", "revenue_per_customer", "eur"),
    ],
    "outlook": [("YoY Revenue Growth", "yoy_revenue_growth_pct", "pct"), ("ROCE", "roce_pct", "pct"), ("DSCR", "dscr", "ratio")],
}


def _format_value(value, kind: str) -> str:
    if value is None:
        return "—"
    value = float(value)
    if kind == "eur":
        return f"€{value:,.0f}"
    if kind == "pct":
        return f"{value:.1f}%"
    if kind == "months":
        return f"{value:.1f} mo"
    if kind == "ratio":
        return f"{value:.2f}x"
    if kind == "count":
        return f"{value:,.0f}"
    return str(value)


def _bullets_from_text(text: str, max_bullets: int = 2) -> list[str]:
    if not text:
        return []
    sentences = [s.strip() for s in text.replace("\n", " ").split(". ") if s.strip()]
    bullets = []
    for s in sentences[:max_bullets]:
        bullets.append(s if s.endswith(".") else f"{s}.")
    return bullets


def _stat_values(spec, section_key: str) -> list[tuple[str, str]]:
    period = spec.period
    fields = STAT_SLIDE_FIELDS.get(section_key, [])
    if section_key == "outlook":
        source = {
            "yoy_revenue_growth_pct": period.yoy_revenue_growth_pct,
            "roce_pct": period.roce_pct,
            "dscr": period.dscr,
        }
    else:
        source = SUMMARY_BUILDERS[section_key](period) or {}

    stats = []
    for label, key, kind in fields:
        if key in source:
            stats.append((label, _format_value(source.get(key), kind)))
    return stats


def _add_textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def _add_title_slide(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    period = spec.period
    _add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1), "Senus PLC", 24, GREY)
    _add_textbox(
        slide, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.2),
        spec.title or f"{period.label} Financial Report", 40, INK, bold=True,
    )
    _add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.6), f"Prepared for {spec.audience_label}", 20, FOREST)
    if spec.context_note:
        _add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1), spec.context_note, 14, GREY, italic=True)


def _add_section_slide(prs, layout, spec, section_key):
    slide = prs.slides.add_slide(layout)
    _add_textbox(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8), SECTION_LABELS[section_key], 28, INK, bold=True)

    top = Inches(1.6)
    for label, value in _stat_values(spec, section_key):
        _add_textbox(slide, Inches(0.6), top, Inches(4), Inches(0.4), label, 14, GREY)
        _add_textbox(slide, Inches(0.6), top + Inches(0.35), Inches(4), Inches(0.6), value, 28, FOREST, bold=True)
        top += Inches(1.1)

    if spec.narrative_approved and spec.tailored_narrative:
        bullets = _bullets_from_text(spec.tailored_narrative.get(section_key, ""))
        bullet_top = Inches(1.6)
        for bullet in bullets:
            _add_textbox(slide, Inches(5.2), bullet_top, Inches(7.3), Inches(1), f"•  {bullet}", 16, INK)
            bullet_top += Inches(1.0)


def generate_deck(spec) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    _add_title_slide(prs, blank_layout, spec)
    for section_key, field_name in spec.SECTION_FIELDS:
        if getattr(spec, field_name):
            _add_section_slide(prs, blank_layout, spec, section_key)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
